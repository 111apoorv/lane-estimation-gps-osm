"""
Create PowerPoint Matching Gamma Presentation
Lane Estimation Using GPS Probe Data for OSM Imputation
By Apoorv Mishra | Rutgers Business School | Fall 2025
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# Helper function for RGB colors
def rgb_color(r, g, b):
    return f"{r:02X}{g:02X}{b:02X}"

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors (matching Gamma dark theme) - as hex strings
DARK_BG = rgb_color(20, 20, 35)
ORANGE = rgb_color(255, 140, 66)
WHITE = rgb_color(255, 255, 255)
LIGHT_GRAY = rgb_color(180, 180, 180)
CARD_BG = rgb_color(35, 35, 55)

def set_shape_color(shape, hex_color):
    """Set fill color of a shape"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = type('obj', (object,), {'__int__': lambda s: int(hex_color, 16)})()
    # Workaround for setting color
    fill = shape.fill._xPr
    solidFill = fill.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            srgbClr.set('val', hex_color)

def set_font_color(paragraph, hex_color):
    """Set font color"""
    for run in paragraph.runs:
        run.font.color.rgb = None
    paragraph.font.color.rgb = None
    # Set via XML
    try:
        from pptx.oxml.ns import nsmap
        rPr = paragraph._p.find(qn('a:r'))
        if rPr is not None:
            rPr_elem = rPr.find(qn('a:rPr'))
            if rPr_elem is not None:
                solidFill = rPr_elem.find(qn('a:solidFill'))
                if solidFill is not None:
                    srgbClr = solidFill.find(qn('a:srgbClr'))
                    if srgbClr is not None:
                        srgbClr.set('val', hex_color)
    except:
        pass

def add_dark_slide(prs):
    """Add a slide with dark background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    return slide

def add_title_text(slide, text, top, size=44, color=WHITE, bold=True):
    """Add title text"""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return box

def add_body_text(slide, text, top, size=20, color=LIGHT_GRAY, left=0.5, width=12.3):
    """Add body text"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return box

def add_card(slide, title, content, left, top, width=3.8, height=2.5):
    """Add a card with title and content"""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = ORANGE
    card.line.width = Pt(2)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(width - 0.4), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.8), Inches(width - 0.4), Inches(height - 1))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_numbered_item(slide, number, title, content, left, top, width=3.8):
    """Add numbered item with card"""
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.05), Inches(0.5), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(left + 0.6), Inches(top), Inches(width - 0.6), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(left + 0.6), Inches(top + 0.5), Inches(width - 0.6), Inches(1))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY

def add_stat_box(slide, value, label, left, top, width=3.5):
    """Add a statistics box"""
    # Background
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(2))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = ORANGE
    box.line.width = Pt(2)
    
    # Value
    val_box = slide.shapes.add_textbox(Inches(left), Inches(top + 0.3), Inches(width), Inches(0.8))
    tf = val_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER
    
    # Label
    lbl_box = slide.shapes.add_textbox(Inches(left), Inches(top + 1.2), Inches(width), Inches(0.6))
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

print("Creating Gamma-style presentation...")

# ============================================
# SLIDE 1: Title Slide
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Lane Estimation Using GPS Probe Data\nfor OSM Imputation", 2.0, size=42)
add_body_text(slide, "Machine Learning Approach for Lane Count Prediction & Map Enhancement", 4.0, size=24, color=ORANGE)
add_body_text(slide, "Apoorv Mishra | Rutgers Business School | Fall 2025", 5.5, size=18, color=LIGHT_GRAY)

# ============================================
# SLIDE 2: The Challenge
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "The Challenge: Bridging the Gap in Lane Data", 0.5, size=36)
add_body_text(slide, "OpenStreetMap (OSM) often lacks comprehensive or accurate lane count information,\nleading to significant gaps in critical map data.", 1.3, size=18, color=LIGHT_GRAY)

add_numbered_item(slide, 1, "Missing Tags", "Many road segments within OSM are missing crucial lane count attributes.", 0.8, 2.5, width=3.5)
add_numbered_item(slide, 2, "Incorrect Data", "Existing lane data can be inaccurate, leading to inconsistencies in mapping.", 5.0, 2.5, width=3.5)
add_numbered_item(slide, 3, "Real-World Impact", "This data gap affects navigation systems, traffic simulations, and urban planning.", 9.2, 2.5, width=3.5)

# ============================================
# SLIDE 3: Project Objective
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Project Objective: Automated Lane Count Prediction", 0.5, size=36)
add_body_text(slide, "Our project aims to leverage machine learning to predict lane counts,\nthereby enhancing the accuracy and completeness of OpenStreetMap data.", 1.3, size=18, color=LIGHT_GRAY)

# Flow diagram
add_card(slide, "1. GPS Probe Traces", "Utilising raw GPS data from vehicles", 0.5, 3.0, width=2.8, height=1.8)
add_card(slide, "2. Feature Extraction", "Transforming traces into meaningful features", 3.6, 3.0, width=2.8, height=1.8)
add_card(slide, "3. Machine Learning", "Applying ML models for prediction", 6.7, 3.0, width=2.8, height=1.8)
add_card(slide, "4. Lane Prediction", "Outputting estimated lane counts", 9.8, 3.0, width=2.8, height=1.8)

add_body_text(slide, "Predictions are validated against official Florida DOT lane counts to ensure ground truth accuracy.", 5.5, size=16, color=LIGHT_GRAY)

# ============================================
# SLIDE 4: Core Insight
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Core Insight: GPS Clustering Reveals Lane Structure", 0.5, size=36)
add_body_text(slide, "The distribution of GPS probe data across a road's width provides\nunique patterns that indicate the number of lanes present.", 1.3, size=18, color=LIGHT_GRAY)

add_card(slide, "Lane-Dependent Clustering", "Vehicles naturally cluster within their respective lanes as they travel.", 0.8, 2.8, width=3.6, height=2.0)
add_card(slide, "Lateral Distribution", "Analysing the lateral spread reveals distinct peaks corresponding to lane positions.", 4.8, 2.8, width=3.6, height=2.0)
add_card(slide, "Increased Complexity", "More lanes result in a wider spread and more identifiable peaks in the data.", 8.8, 2.8, width=3.6, height=2.0)

# Quote
quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(10), Inches(0.8))
tf = quote_box.text_frame
p = tf.paragraphs[0]
p.text = '"Cars leave quantifiable digital footprints that reveal road infrastructure."'
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 5: Datasets Overview
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Datasets Overview: Fueling the Prediction Model", 0.5, size=36)
add_body_text(slide, "Our project leverages three distinct datasets, each playing a crucial role\nin training, validation, and feature enhancement.", 1.3, size=18, color=LIGHT_GRAY)

add_numbered_item(slide, 1, "LA GPS Probe (Training)", "329,787 filtered road segments provide the core data for model training.", 0.8, 2.8, width=3.8)
add_numbered_item(slide, 2, "Florida DOT Lanes (Validation)", "86,880 segments serve as ground truth for robust validation.", 4.8, 2.8, width=3.8)
add_numbered_item(slide, 3, "Florida DOT AADT (Feature)", "20,289 segments contribute traffic volume data, enriching features.", 8.8, 2.8, width=3.8)

add_body_text(slide, "Each dataset is instrumental in developing a robust and accurate lane estimation model.", 5.5, size=16, color=LIGHT_GRAY)

# ============================================
# SLIDE 6: Dataset 1 Details
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Dataset 1 Details: LA GPS Probe Data", 0.5, size=36)
add_body_text(slide, "The primary training dataset consists of crowdsourced GPS probe traces,\nmeticulously filtered to ensure data quality and relevance.", 1.2, size=18, color=LIGHT_GRAY)

details = """• Source: Crowdsourced GPS probe traces (smartphones, fleet vehicles, navigation systems)

• File: la_link_dist_vs_lanes_estimation_dataset_v1.parquet

• Filtering Criteria:
    - total_count > 1,000 (ensuring sufficient data points per segment)
    - lanes_int ≤ 7 (focusing on common road configurations)

• Target Variable: lanes_int (the integer representation of lane count)

• Final Size: 329,787 road segments after filtering"""

detail_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(4))
tf = detail_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = details
p.font.size = Pt(18)
p.font.color.rgb = WHITE

# ============================================
# SLIDE 7: Feature Engineering
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Feature Engineering: Unlocking Lane Patterns", 0.5, size=36)
add_body_text(slide, "To capture the lateral distribution of GPS points, we devised a robust\nfeature engineering strategy using 40 distinct bins across the road width.", 1.2, size=18, color=LIGHT_GRAY)

add_card(slide, "Road Width Discretisation", "The estimated road width is divided into 40 equal slices (bins).", 0.8, 2.8, width=3.6, height=2.2)
add_card(slide, "GPS Point Ratio", "Each bin represents the ratio of GPS points falling within that slice, summing to 1.", 4.8, 2.8, width=3.6, height=2.2)
add_card(slide, "Peak Identification", "Peaks within this distribution directly indicate the presence of lanes.", 8.8, 2.8, width=3.6, height=2.2)

add_body_text(slide, "An additional 'oneway' flag is incorporated, bringing total features to 41.", 5.8, size=16, color=ORANGE)

# ============================================
# SLIDE 8: Model Comparison
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Model Comparison: Identifying the Best Performer", 0.5, size=36)
add_body_text(slide, "We evaluated several machine learning models to determine\nthe most effective approach for lane count prediction.", 1.2, size=18, color=LIGHT_GRAY)

# Model comparison table-like cards
add_stat_box(slide, "90.11%", "Random Forest\n(WINNER)", 0.8, 2.8, width=2.8)
add_stat_box(slide, "85.0%", "XGBoost", 3.9, 2.8, width=2.8)
add_stat_box(slide, "84.56%", "LightGBM", 7.0, 2.8, width=2.8)
add_stat_box(slide, "22.0%", "GMM\n(Unsupervised)", 10.1, 2.8, width=2.8)

add_body_text(slide, "Random Forest emerged as the top performer due to its ability to handle\nnon-linear relationships and its inherent robustness.", 5.5, size=16, color=ORANGE)

# ============================================
# SLIDE 9: Random Forest Configuration
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Random Forest: Configuration for Optimal Results", 0.5, size=36)
add_body_text(slide, "The Random Forest model was meticulously configured\nto achieve maximum predictive accuracy in lane estimation.", 1.2, size=18, color=LIGHT_GRAY)

add_card(slide, "Number of Estimators", "n_estimators = 200\nEnsures a robust ensemble of decision trees", 0.8, 2.6, width=5.8, height=1.8)
add_card(slide, "Class Weighting", "class_weight = balanced_subsample\nAddresses potential class imbalance", 7.0, 2.6, width=5.8, height=1.8)
add_card(slide, "Feature Set", "40 bin features + oneway flag\nTotal: 41 features", 0.8, 4.8, width=5.8, height=1.8)
add_card(slide, "Train/Test Split", "80/20 stratified split by lane count\nEnsures representative data", 7.0, 4.8, width=5.8, height=1.8)

# ============================================
# SLIDE 10: Main Results
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Main Results: High Accuracy in Lane Prediction", 0.5, size=36)
add_body_text(slide, "Our Random Forest model achieved exceptional accuracy,\nparticularly for predictions within a crucial ±1 lane margin.", 1.2, size=18, color=LIGHT_GRAY)

add_stat_box(slide, "90.11%", "Overall Accuracy\nDirect match to ground truth", 2.5, 2.8, width=4)
add_stat_box(slide, "99.0%", "Within ±1 Lane\nHighly reliable for applications", 7.0, 2.8, width=4)

add_body_text(slide, "A substantial majority of predictions are either perfectly correct or off by only a single lane,\ndemonstrating the model's practical utility for OSM data imputation.", 5.5, size=16, color=LIGHT_GRAY)

# ============================================
# SLIDE 11: Error Analysis
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Error Analysis: Understanding Model Performance", 0.5, size=36)

add_stat_box(slide, "59,438", "Correct Predictions\n(90.1%)", 0.5, 1.8, width=3)
add_stat_box(slide, "5,858", "Off by 1 Lane\n(8.9%)", 3.8, 1.8, width=3)
add_stat_box(slide, "655", "Off by 2 Lanes\n(1.0%)", 7.1, 1.8, width=3)
add_stat_box(slide, "7", "Off by 3+ Lanes\n(0.0%)", 10.4, 1.8, width=2.5)

# Key findings
findings = """Key Findings:
• 1-lane and 3-lane roads show highest error rates (27%)
• 2-lane roads are predicted most accurately (2.8% error)
• Most common mistake: 3 lanes predicted as 2 lanes (1,944 times)
• Adjacent lane counts are easily confused due to similar GPS patterns"""

findings_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(12), Inches(2.5))
tf = findings_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = findings
p.font.size = Pt(18)
p.font.color.rgb = WHITE

# ============================================
# SLIDE 12: AADT Discovery
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Key Discovery: AADT Correlation with Lane Count", 0.5, size=36)
add_body_text(slide, "Analysis of Florida DOT data revealed a strong correlation\nbetween traffic volume (AADT) and lane count.", 1.2, size=18, color=LIGHT_GRAY)

# AADT values
aadt_data = """Traffic Volume by Lane Count:

1 lane  →  13,249 vehicles/day
2 lanes →  20,390 vehicles/day
3 lanes →  52,711 vehicles/day
4 lanes →  79,403 vehicles/day
5 lanes →  142,727 vehicles/day
6 lanes →  149,307 vehicles/day

INSIGHT: More Traffic = More Lanes"""

aadt_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(6), Inches(4))
tf = aadt_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = aadt_data
p.font.size = Pt(18)
p.font.color.rgb = WHITE

add_stat_box(slide, "43%", "AADT-Only Accuracy\n(Just 1 feature!)", 8, 2.8, width=4)
add_stat_box(slide, "87.7%", "Within 1 Lane\nUsing AADT alone", 8, 5.2, width=4)

# ============================================
# SLIDE 13: Proposed Improvement
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Proposed Improvement: Adding AADT Feature", 0.5, size=36)

# Current vs Enhanced
add_card(slide, "Current Model", "Features: 41\n(40 GPS bins + oneway)\n\nAccuracy: 90.11%\n\nLimitation: 2↔3 lane confusion", 1, 1.8, width=5.5, height=3.5)
add_card(slide, "Enhanced Model", "Features: 42\n(40 GPS + oneway + AADT)\n\nExpected: 93-95%\n\nAADT distinguishes 2↔3 lanes", 7, 1.8, width=5.5, height=3.5)

add_body_text(slide, "Adding AADT helps distinguish 2-lane vs 3-lane roads\n(20K vs 53K vehicles/day - a clear signal!)", 5.8, size=18, color=ORANGE)

# ============================================
# SLIDE 14: Conclusion
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Conclusion & Future Work", 0.5, size=36)

conclusion = """Key Achievements:
✓ 90.11% accuracy using GPS probe distribution features
✓ 99% of predictions within ±1 lane of true value
✓ Identified AADT as key feature for improvement (→ 93-95%)
✓ Validated against Florida DOT ground truth (86,880 roads)

Future Work:
• Integrate AADT feature into production model
• Expand to additional states and regions
• Deploy for large-scale OSM lane data imputation
• Explore additional DOT datasets nationwide"""

concl_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
tf = concl_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = conclusion
p.font.size = Pt(22)
p.font.color.rgb = WHITE

# ============================================
# SLIDE 15: Thank You
# ============================================
slide = add_dark_slide(prs)
add_title_text(slide, "Thank You!", 2.5, size=52)
add_body_text(slide, "Questions?", 4.0, size=32, color=ORANGE)
add_body_text(slide, "Apoorv Mishra | Rutgers Business School", 5.5, size=18, color=LIGHT_GRAY)

# Save presentation
output_path = r"C:\Users\webap\Downloads\Lane identification\Lane_Estimation_Presentation_Gamma_Style.pptx"
prs.save(output_path)

print(f"\n✅ Presentation created successfully!")
print(f"📁 Saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")

